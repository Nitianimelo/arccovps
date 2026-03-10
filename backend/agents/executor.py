"""
Execução de ferramentas para os agentes especialistas.

Cada função de ferramenta está isolada aqui, desacoplada do endpoint HTTP.
"""

import asyncio
import io
import logging
import os
import tempfile
import time

import httpx

logger = logging.getLogger(__name__)


async def execute_tool(func_name: str, func_args: dict) -> str:
    """Despachante principal: executa a ferramenta e retorna resultado como string."""
    if func_name == "web_search":
        return await _web_search(func_args.get("query", ""))

    elif func_name == "web_fetch":
        return await _web_fetch(func_args.get("url", ""))

    elif func_name == "generate_pdf":
        return await _generate_pdf(func_args)

    elif func_name == "generate_excel":
        return await _generate_excel(func_args)

    elif func_name == "execute_python":
        return await _execute_python(func_args.get("code", ""))

    elif func_name == "fetch_file_content":
        return await _fetch_file_content(func_args.get("url", ""))

    elif func_name == "modify_excel":
        return await _modify_excel(func_args)

    elif func_name == "modify_pptx":
        return await _modify_pptx(func_args)

    elif func_name == "modify_pdf":
        return await _modify_pdf(func_args)

    elif func_name == "ask_browser":
        return await _ask_browser(func_args)

    elif func_name == "generate_pdf_template":
        return await _generate_pdf_template(func_args)

    elif func_name == "use_design_template":
        return await _use_design_template(func_args)

    return f"Ferramenta desconhecida: {func_name}"


# ── Implementações ─────────────────────────────────────────────────────────────

async def _web_search(query: str) -> str:
    from backend.services.search_service import search_web_formatted
    return await search_web_formatted(query)


async def _web_fetch(url: str) -> str:
    try:
        from bs4 import BeautifulSoup
        async with httpx.AsyncClient(timeout=20.0) as client:
            response = await client.get(
                url, headers={"User-Agent": "ArccoAgent/2.0"}, follow_redirects=True
            )
            html = response.text

        soup = BeautifulSoup(html, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header", "aside", "form", "svg", "noscript"]):
            tag.decompose()

        text = soup.get_text(separator=" ", strip=True)
        title = soup.title.string if soup.title else url

        if len(text) > 20000:
            text = text[:20000] + "... [Truncado]"

        return f"**Conteúdo de {url}**\n**Título:** {title}\n\n{text}"
    except Exception as e:
        return f"Erro ao ler URL ({url}): {e}"


async def _ask_browser(args: dict) -> str:
    """
    Navega remotamente via Browserbase + Playwright (CDP).
    Delega toda a lógica para browser_service.execute_browserbase_task.
    """
    from backend.services.browser_service import execute_browserbase_task

    url = args.get("url", "")
    if not url:
        return "Erro: URL não fornecida para o Browser Agent."

    return await execute_browserbase_task(
        url=url,
        actions=args.get("actions", []),
        wait_for=int(args.get("wait_for") or 0),
        mobile=bool(args.get("mobile", False)),
        include_tags=args.get("include_tags"),
        exclude_tags=args.get("exclude_tags"),
    )


async def _generate_pdf(args: dict) -> str:
    from backend.core.config import get_config
    from backend.core.supabase_client import upload_to_supabase

    config = get_config()
    html_content = (args.get("html_content") or "").strip()

    if html_content:
        # Modo Playwright: HTML+Tailwind → PDF de alta qualidade
        from backend.services.file_service import generate_pdf_playwright
        logger.info("[PDF] Modo Playwright (HTML+Tailwind)")
        pdf_bytes = await generate_pdf_playwright(html_content)
    else:
        # Modo texto: reportlab (fallback)
        from backend.services.file_service import generate_pdf
        logger.info("[PDF] Modo reportlab (texto)")
        title = args.get("title", "documento")
        content = args.get("content", "")
        pdf_bytes = await asyncio.to_thread(generate_pdf, title, content)

    filename = args.get("filename", f"doc-{int(time.time())}")
    if not filename.endswith(".pdf"):
        filename += ".pdf"

    url = await asyncio.to_thread(
        upload_to_supabase,
        config.supabase_storage_bucket,
        filename,
        pdf_bytes,
        "application/pdf",
    )
    return (
        f"PDF gerado com sucesso. URL: {url}\n\n"
        f"INSTRUÇÃO OBRIGATÓRIA: Inclua exatamente este link na resposta final: [Baixar PDF]({url})"
    )


async def _generate_pdf_template(args: dict) -> str:
    from backend.core.config import get_config
    from backend.core.supabase_client import upload_to_supabase
    from backend.services.file_service import generate_pdf_from_template

    config = get_config()
    template_name = args.get("template_name", "relatorio")
    data = args.get("data", {})

    try:
        logger.info(f"[PDF] Gerando via template Jinja2: '{template_name}'")
        pdf_bytes = await generate_pdf_from_template(template_name, data)
    except Exception as e:
        return f"Erro ao gerar PDF com template '{template_name}': {e}"

    filename = args.get("filename", f"{template_name}-{int(time.time())}")
    if not filename.endswith(".pdf"):
        filename += ".pdf"

    url = await asyncio.to_thread(
        upload_to_supabase,
        config.supabase_storage_bucket,
        filename,
        pdf_bytes,
        "application/pdf",
    )
    return (
        f"PDF gerado via template '{template_name}'. URL: {url}\n\n"
        f"INSTRUÇÃO OBRIGATÓRIA: Inclua exatamente este link na resposta final: [Baixar PDF]({url})"
    )


async def _use_design_template(args: dict) -> str:
    """Aplica um template de design pré-construído e retorna HTML standalone."""
    from backend.services.template_service import get_template_html, apply_content, search_pexels_image

    slug = (args.get("slug") or "").strip()
    if not slug:
        return "Erro: campo 'slug' obrigatório. Verifique o catálogo no system prompt."

    try:
        html = get_template_html(slug)
    except FileNotFoundError as e:
        return str(e)
    except Exception as e:
        return f"Erro ao ler template '{slug}': {e}"

    content = {
        k: args[k]
        for k in ("title", "eyebrow", "subtitle", "footer", "heading", "extra_patches")
        if args.get(k) is not None
    }

    # Resolve imagem: pexels_query tem prioridade sobre image_url direto
    image_url = args.get("image_url") or ""
    pexels_query = (args.get("pexels_query") or "").strip()
    if not image_url and pexels_query:
        image_url = await search_pexels_image(pexels_query)

    try:
        html = apply_content(
            html,
            content,
            color_overrides=args.get("color_overrides"),
            image_url=image_url or None,
        )
    except Exception as e:
        logger.warning(f"[TEMPLATE] Erro ao aplicar conteúdo no template '{slug}': {e}")
        # Retorna o template mesmo sem patches — melhor do que falhar

    return html


async def _generate_excel(args: dict) -> str:
    from backend.core.config import get_config
    from backend.core.supabase_client import upload_to_supabase
    from openpyxl import Workbook

    config = get_config()

    def sync_excel():
        headers = [str(h) for h in args.get("headers", [])]
        rows = [[str(c) for c in row] for row in args.get("rows", [])]
        title = args.get("title", "Planilha")[:31]

        wb = Workbook()
        ws = wb.active
        ws.title = title
        ws.append(headers)
        for row in rows:
            ws.append(row)

        buffer = io.BytesIO()
        wb.save(buffer)
        file_bytes = buffer.getvalue()

        filename = args.get("filename", f"planilha-{int(time.time())}")
        if not filename.endswith(".xlsx"):
            filename += ".xlsx"

        return upload_to_supabase(
            config.supabase_storage_bucket,
            filename,
            file_bytes,
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )

    url = await asyncio.to_thread(sync_excel)
    return (
        f"Planilha Excel gerada. URL: {url}\n\n"
        f"INSTRUÇÃO OBRIGATÓRIA: Inclua exatamente este link na resposta final: [Baixar Planilha]({url})"
    )


async def _execute_python(code: str) -> str:
    from backend.core.config import get_config
    config = get_config()

    if not config.allow_code_execution:
        return "❌ Execução de código desabilitada neste ambiente."

    blocked = [
        "os.system", "os.popen", "os.exec", "os.spawn", "os.remove", "os.unlink", "os.rmdir",
        "eval(", "exec(", "__import__", "requests.", "urllib.",
        "subprocess", "importlib", "shutil.rmtree", "shutil.move",
        "socket.", "ctypes", "getattr(", "setattr(", "delattr(",
        "open(", "pathlib",
    ]
    for b in blocked:
        if b in code:
            return f"❌ Operação bloqueada por segurança: {b}"

    tmp_dir = tempfile.gettempdir()
    with tempfile.NamedTemporaryFile(mode="w", suffix=".py", delete=False, dir=tmp_dir) as f:
        f.write(code)
        tmp_name = f.name

    try:
        python_cmd = "python" if os.name == "nt" else "python3"
        process = await asyncio.create_subprocess_exec(
            python_cmd, tmp_name,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
            cwd=tmp_dir,
        )
        try:
            stdout, stderr = await asyncio.wait_for(process.communicate(), timeout=10.0)
        except asyncio.TimeoutError:
            process.kill()
            await process.communicate()
            return "❌ Timeout na execução (10s excedidos)."

        out = stdout.decode() if stdout else ""
        err = stderr.decode() if stderr else ""
        result = out + (f"\nSTDERR: {err}" if err else "")
        return result if result.strip() else "(Código executado sem output. Use print() para ver resultados.)"
    except Exception as e:
        return f"Erro na execução: {e}"
    finally:
        if os.path.exists(tmp_name):
            os.unlink(tmp_name)


# ── Modificador de Arquivos ────────────────────────────────────────────────────

async def _fetch_file_content(url: str) -> str:
    """Baixa um arquivo e retorna sua estrutura como texto legível."""
    try:
        async with httpx.AsyncClient(timeout=30.0) as client:
            response = await client.get(url, follow_redirects=True)
        if response.status_code != 200:
            return f"Erro ao baixar arquivo: HTTP {response.status_code}"

        content_type = response.headers.get("content-type", "").lower()
        file_bytes = response.content
        url_lower = url.lower().split("?")[0]  # ignora query params

        if "spreadsheet" in content_type or url_lower.endswith(".xlsx"):
            return await asyncio.to_thread(_read_excel_structure, file_bytes)
        elif "presentation" in content_type or url_lower.endswith(".pptx"):
            return await asyncio.to_thread(_read_pptx_structure, file_bytes)
        elif "pdf" in content_type or url_lower.endswith(".pdf"):
            return await asyncio.to_thread(_read_pdf_text, file_bytes)
        else:
            return f"Tipo de arquivo não identificado (content-type: {content_type}). URL: {url}"
    except Exception as e:
        return f"Erro ao ler arquivo: {e}"


def _read_excel_structure(file_bytes: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(file_bytes))
    lines = [f"Planilha Excel — {len(wb.sheetnames)} aba(s): {', '.join(wb.sheetnames)}"]
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            lines.append(f"\nAba '{sheet_name}': vazia")
            continue
        headers = [str(c) if c is not None else "" for c in rows[0]]
        lines.append(f"\nAba '{sheet_name}' — {ws.max_row} linha(s), {ws.max_column} coluna(s)")
        lines.append(f"Cabeçalhos (linha 1): {headers}")
        for i, row in enumerate(rows[1:6], start=2):
            lines.append(f"  Linha {i}: {[str(c) if c is not None else '' for c in row]}")
        if ws.max_row > 6:
            lines.append(f"  ... ({ws.max_row - 6} linha(s) restante(s) não exibidas)")
    return "\n".join(lines)


def _read_pptx_structure(file_bytes: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    lines = [f"Apresentação PPTX — {len(prs.slides)} slide(s)"]
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
        preview = " | ".join(texts[:4]) if texts else "(sem texto)"
        lines.append(f"  Slide {i + 1}: {preview}")
    return "\n".join(lines)


def _read_pdf_text(file_bytes: bytes) -> str:
    import PyPDF2
    reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
    lines = [f"PDF — {len(reader.pages)} página(s)"]
    total_chars = 0
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        snippet = text[:800].strip()
        lines.append(f"\n--- Página {i + 1} ---\n{snippet}")
        total_chars += len(text)
        if total_chars > 4000:
            lines.append("... [restante truncado]")
            break
    return "\n".join(lines)


async def _modify_excel(args: dict) -> str:
    from backend.core.config import get_config
    from backend.core.supabase_client import upload_to_supabase

    config = get_config()
    url = args.get("url", "")
    cell_updates = args.get("cell_updates", [])
    append_rows = args.get("append_rows", [])
    output_filename = args.get("output_filename", f"planilha-modificada")

    try:
        async with httpx.AsyncClient(timeout=30.0) as client:
            response = await client.get(url, follow_redirects=True)
        if response.status_code != 200:
            return f"Erro ao baixar planilha: HTTP {response.status_code}"
        file_bytes = response.content
    except Exception as e:
        return f"Erro ao baixar planilha: {e}"

    def sync_modify():
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(file_bytes))

        for update in cell_updates:
            sheet_name = update.get("sheet", "")
            ws = wb[sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.active
            ws[update["cell"]] = update["value"]

        for row_def in append_rows:
            sheet_name = row_def.get("sheet", "")
            ws = wb[sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.active
            ws.append(row_def.get("values", []))

        buffer = io.BytesIO()
        wb.save(buffer)
        return buffer.getvalue()

    modified_bytes = await asyncio.to_thread(sync_modify)
    filename = output_filename if output_filename.endswith(".xlsx") else output_filename + ".xlsx"
    upload_url = await asyncio.to_thread(
        upload_to_supabase,
        config.supabase_storage_bucket,
        filename,
        modified_bytes,
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )
    return (
        f"Planilha modificada com sucesso. URL: {upload_url}\n\n"
        f"INSTRUÇÃO OBRIGATÓRIA: Inclua exatamente este link na resposta final: [Baixar Planilha Modificada]({upload_url})"
    )


async def _modify_pptx(args: dict) -> str:
    from backend.core.config import get_config
    from backend.core.supabase_client import upload_to_supabase

    config = get_config()
    url = args.get("url", "")
    text_replacements = args.get("text_replacements", [])
    output_filename = args.get("output_filename", f"apresentacao-modificada")

    try:
        async with httpx.AsyncClient(timeout=30.0) as client:
            response = await client.get(url, follow_redirects=True)
        if response.status_code != 200:
            return f"Erro ao baixar apresentação: HTTP {response.status_code}"
        file_bytes = response.content
    except Exception as e:
        return f"Erro ao baixar apresentação: {e}"

    def sync_modify():
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        for rep in text_replacements:
                            if rep["find"] in run.text:
                                run.text = run.text.replace(rep["find"], rep["replace"])

        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

    modified_bytes = await asyncio.to_thread(sync_modify)
    filename = output_filename if output_filename.endswith(".pptx") else output_filename + ".pptx"
    upload_url = await asyncio.to_thread(
        upload_to_supabase,
        config.supabase_storage_bucket,
        filename,
        modified_bytes,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    return (
        f"Apresentação modificada com sucesso. URL: {upload_url}\n\n"
        f"INSTRUÇÃO OBRIGATÓRIA: Inclua exatamente este link na resposta final: [Baixar Apresentação Modificada]({upload_url})"
    )


async def _modify_pdf(args: dict) -> str:
    from backend.core.config import get_config
    from backend.core.supabase_client import upload_to_supabase

    config = get_config()
    url = args.get("url", "")
    text_replacements = args.get("text_replacements", [])
    append_content = args.get("append_content", "")
    output_filename = args.get("output_filename", f"documento-modificado")

    try:
        async with httpx.AsyncClient(timeout=30.0) as client:
            response = await client.get(url, follow_redirects=True)
        if response.status_code != 200:
            return f"Erro ao baixar PDF: HTTP {response.status_code}"
        file_bytes = response.content
    except Exception as e:
        return f"Erro ao baixar PDF: {e}"

    def sync_modify():
        import PyPDF2
        from reportlab.lib.pagesizes import A4
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet

        # Extrai texto de todas as páginas
        reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
        full_text = "\n\n".join(page.extract_text() or "" for page in reader.pages)

        # Aplica substituições
        for rep in text_replacements:
            full_text = full_text.replace(rep["find"], rep["replace"])

        # Adiciona conteúdo extra
        if append_content:
            full_text += f"\n\n{append_content}"

        # Regera o PDF com reportlab
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []
        for line in full_text.split("\n"):
            if line.strip():
                safe_line = line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                story.append(Paragraph(safe_line, styles["Normal"]))
                story.append(Spacer(1, 4))
        doc.build(story)
        return buffer.getvalue()

    modified_bytes = await asyncio.to_thread(sync_modify)
    filename = output_filename if output_filename.endswith(".pdf") else output_filename + ".pdf"
    upload_url = await asyncio.to_thread(
        upload_to_supabase,
        config.supabase_storage_bucket,
        filename,
        modified_bytes,
        "application/pdf",
    )
    return (
        f"PDF modificado com sucesso. URL: {upload_url}\n\n"
        f"INSTRUÇÃO OBRIGATÓRIA: Inclua exatamente este link na resposta final: [Baixar PDF Modificado]({upload_url})"
    )
