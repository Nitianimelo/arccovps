"""
Serviço de geração de arquivos — PDF, DOCX, XLSX, PPTX.
Portado de netlify/functions/files.ts
"""

import asyncio
import io
import json
import logging
import os
from pathlib import Path
from typing import Tuple

logger = logging.getLogger(__name__)

_TEMPLATES_DIR = Path(__file__).parent / "pdf_templates"


async def generate_pdf_playwright(html_content: str) -> bytes:
    """
    Gera PDF de alta qualidade usando Playwright para renderizar HTML+Tailwind CSS.
    Injeta automaticamente o CDN do Tailwind se não estiver presente.
    """

    def _sync_render() -> bytes:
        from playwright.sync_api import sync_playwright

        # Injeta Tailwind CDN caso o HTML não tenha estilo próprio
        tailwind_cdn = '<script src="https://cdn.tailwindcss.com"></script>'
        if "tailwindcss" not in html_content and "cdn.tailwindcss" not in html_content:
            if "<head>" in html_content:
                inject_html = html_content.replace("<head>", f"<head>{tailwind_cdn}", 1)
            elif "<html" in html_content:
                inject_html = html_content.replace("<html", f"{tailwind_cdn}<html", 1)
            else:
                inject_html = tailwind_cdn + html_content
        else:
            inject_html = html_content

        with sync_playwright() as p:
            browser = p.chromium.launch(args=["--no-sandbox", "--disable-dev-shm-usage"])
            page = browser.new_page()
            try:
                page.set_content(inject_html, wait_until="networkidle", timeout=30_000)
                pdf_bytes = page.pdf(
                    format="A4",
                    print_background=True,
                    margin={"top": "1.5cm", "right": "1.5cm", "bottom": "1.5cm", "left": "1.5cm"},
                )
            finally:
                browser.close()
        return pdf_bytes

    return await asyncio.to_thread(_sync_render)


async def generate_pdf_from_template(template_name: str, data: dict) -> bytes:
    """
    Gera PDF a partir de um template Jinja2 HTML pré-aprovado.
    O LLM fornece apenas os dados (JSON); o layout vem do template.
    """
    try:
        from jinja2 import Environment, FileSystemLoader, select_autoescape
    except ImportError:
        raise RuntimeError("Jinja2 não instalado. Execute: pip install jinja2")

    template_file = f"{template_name}.html"
    if not (_TEMPLATES_DIR / template_file).exists():
        available = [p.stem for p in _TEMPLATES_DIR.glob("*.html")]
        raise ValueError(
            f"Template '{template_name}' não encontrado. Disponíveis: {available}"
        )

    env = Environment(
        loader=FileSystemLoader(str(_TEMPLATES_DIR)),
        autoescape=select_autoescape(["html"]),
    )
    template = env.get_template(template_file)
    html_content = template.render(**data)
    return await generate_pdf_playwright(html_content)


def generate_pdf(title: str, content: str) -> bytes:
    """Gera PDF com reportlab."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.units import cm

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4,
                            topMargin=2*cm, bottomMargin=2*cm,
                            leftMargin=2*cm, rightMargin=2*cm)
    styles = getSampleStyleSheet()

    title_style = ParagraphStyle(
        "CustomTitle", parent=styles["Title"],
        fontSize=18, spaceAfter=20
    )
    body_style = ParagraphStyle(
        "CustomBody", parent=styles["Normal"],
        fontSize=12, leading=16
    )

    elements = [Paragraph(title, title_style), Spacer(1, 12)]

    for line in content.split("\n"):
        line = line.strip()
        if not line:
            elements.append(Spacer(1, 8))
        elif line.startswith("# "):
            elements.append(Paragraph(f"<b>{line[2:]}</b>", styles["Heading1"]))
        elif line.startswith("## "):
            elements.append(Paragraph(f"<b>{line[3:]}</b>", styles["Heading2"]))
        else:
            # Escape XML chars for reportlab
            safe = line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            elements.append(Paragraph(safe, body_style))

    doc.build(elements)
    return buffer.getvalue()


def generate_docx(title: str, content: str) -> bytes:
    """Gera DOCX com python-docx."""
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    doc.add_heading(title, level=0)
    doc.add_paragraph("")  # spacer

    for line in content.split("\n"):
        trimmed = line.strip()
        if trimmed.startswith("# "):
            doc.add_heading(trimmed[2:], level=1)
        elif trimmed.startswith("## "):
            doc.add_heading(trimmed[3:], level=2)
        elif trimmed.startswith("### "):
            doc.add_heading(trimmed[4:], level=3)
        elif trimmed:
            p = doc.add_paragraph(trimmed)
            for run in p.runs:
                run.font.size = Pt(12)

    buffer = io.BytesIO()
    doc.save(buffer)
    return buffer.getvalue()


def generate_xlsx(title: str, content: str) -> bytes:
    """Gera Excel com openpyxl."""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = title[:31] if title else "Dados"

    # Tentar parsear content como JSON (array de objetos ou array de arrays)
    try:
        data = json.loads(content)
        if isinstance(data, list) and len(data) > 0:
            if isinstance(data[0], dict):
                # Array de objetos: headers das keys
                headers = list(data[0].keys())
                ws.append(headers)
                for row in data:
                    ws.append([str(row.get(h, "")) for h in headers])
            elif isinstance(data[0], list):
                # Array de arrays
                for row in data:
                    ws.append([str(cell) for cell in row])
            else:
                ws.append([title])
                ws.append([content])
        else:
            ws.append([title])
            ws.append([content])
    except (json.JSONDecodeError, TypeError):
        # Não é JSON, tratar como texto
        ws.append([title])
        for line in content.split("\n"):
            if line.strip():
                ws.append([line.strip()])

    buffer = io.BytesIO()
    wb.save(buffer)
    return buffer.getvalue()


def generate_pptx(title: str, content: str) -> bytes:
    """Gera PPTX com python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()

    # Separar por marcador SLIDE:
    import re
    slides_content = re.split(r'SLIDE:', content, flags=re.IGNORECASE)
    slides_content = [s.strip() for s in slides_content if s.strip()]

    if not slides_content:
        # Fallback: um slide com tudo
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = content
    else:
        for slide_text in slides_content:
            lines = slide_text.split("\n")
            slide_title = lines[0].replace("*", "").replace("#", "").strip()
            slide_body = "\n".join(lines[1:]).strip()

            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = slide_title
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = slide_body

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def _inject_tailwind_if_needed(html_content: str) -> str:
    """Injeta Tailwind CDN se o HTML não tiver estilos próprios."""
    tailwind_cdn = '<script src="https://cdn.tailwindcss.com"></script>'
    if "tailwindcss" not in html_content and "cdn.tailwindcss" not in html_content:
        if "<head>" in html_content:
            return html_content.replace("<head>", f"<head>{tailwind_cdn}", 1)
        elif "<html" in html_content:
            return tailwind_cdn + html_content
        return tailwind_cdn + html_content
    return html_content


async def html_to_screenshot(html_content: str, img_format: str = "png") -> bytes:
    """
    Captura screenshot de um HTML via Playwright.
    img_format: "png" ou "jpeg"
    """
    inject = _inject_tailwind_if_needed(html_content)
    fmt = img_format.lower() if img_format.lower() in ("png", "jpeg") else "png"

    def _sync() -> bytes:
        from playwright.sync_api import sync_playwright
        with sync_playwright() as p:
            browser = p.chromium.launch(args=["--no-sandbox", "--disable-dev-shm-usage"])
            page = browser.new_page(viewport={"width": 1280, "height": 720})
            try:
                page.set_content(inject, wait_until="networkidle", timeout=30_000)
                data = page.screenshot(full_page=False, type=fmt)
            finally:
                browser.close()
        return data

    return await asyncio.to_thread(_sync)


async def html_to_pptx(html_content: str, title: str = "Apresentação") -> bytes:
    """
    Converte HTML com slides (.slide) em PPTX via screenshots Playwright.
    Cada <section class="slide"> vira um slide independente.
    Se não houver .slide, usa um único slide com screenshot full.
    """
    inject = _inject_tailwind_if_needed(html_content)

    def _sync() -> bytes:
        import io as _io
        from playwright.sync_api import sync_playwright
        from pptx import Presentation
        from pptx.util import Inches

        screenshots: list[bytes] = []

        with sync_playwright() as p:
            browser = p.chromium.launch(args=["--no-sandbox", "--disable-dev-shm-usage"])
            page = browser.new_page(viewport={"width": 1280, "height": 720})
            try:
                page.set_content(inject, wait_until="networkidle", timeout=30_000)
                slide_count = page.evaluate("() => document.querySelectorAll('.slide').length")

                if slide_count and slide_count > 0:
                    for i in range(int(slide_count)):
                        # Mostra apenas o slide i, esconde os outros
                        page.evaluate(f"""() => {{
                            const slides = document.querySelectorAll('.slide');
                            slides.forEach((s, idx) => s.style.display = idx === {i} ? 'flex' : 'none');
                        }}""")
                        page.wait_for_timeout(300)
                        screenshots.append(page.screenshot(full_page=False, type="png"))
                else:
                    screenshots.append(page.screenshot(full_page=False, type="png"))
            finally:
                browser.close()

        # Monta o PPTX
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]

        for img_bytes in screenshots:
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                _io.BytesIO(img_bytes), 0, 0,
                prs.slide_width, prs.slide_height
            )

        buf = _io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    return await asyncio.to_thread(_sync)


def _text_to_html(title: str, content: str) -> str:
    """Converte texto/markdown simples em HTML bonito para exportação via Playwright."""
    import html as html_lib
    lines = content.split("\n")
    body_parts = []
    for line in lines:
        stripped = line.strip()
        if not stripped:
            body_parts.append("<br/>")
        elif stripped.startswith("### "):
            body_parts.append(f'<h3 class="text-lg font-semibold text-gray-700 mt-4 mb-1">{html_lib.escape(stripped[4:])}</h3>')
        elif stripped.startswith("## "):
            body_parts.append(f'<h2 class="text-xl font-bold text-gray-800 mt-6 mb-2">{html_lib.escape(stripped[3:])}</h2>')
        elif stripped.startswith("# "):
            body_parts.append(f'<h1 class="text-2xl font-bold text-gray-900 mt-6 mb-3">{html_lib.escape(stripped[2:])}</h1>')
        elif stripped.startswith("- ") or stripped.startswith("* "):
            body_parts.append(f'<li class="ml-5 list-disc text-gray-700">{html_lib.escape(stripped[2:])}</li>')
        else:
            body_parts.append(f'<p class="text-gray-700 leading-relaxed my-1">{html_lib.escape(stripped)}</p>')

    body_html = "\n".join(body_parts)
    safe_title = html_lib.escape(title)
    return f"""<!DOCTYPE html>
<html lang="pt-BR">
<head>
  <meta charset="UTF-8"/>
  <script src="https://cdn.tailwindcss.com"></script>
  <title>{safe_title}</title>
</head>
<body class="bg-white font-sans p-12 max-w-3xl mx-auto">
  <h1 class="text-3xl font-bold text-gray-900 border-b-2 border-indigo-600 pb-4 mb-8">{safe_title}</h1>
  <div class="prose text-base">
    {body_html}
  </div>
</body>
</html>"""


# Mapa de tipos para geração
FILE_GENERATORS = {
    "pdf": {
        "func": generate_pdf,
        "mime": "application/pdf",
        "ext": "pdf",
    },
    "docx": {
        "func": generate_docx,
        "mime": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "ext": "docx",
    },
    "excel": {
        "func": generate_xlsx,
        "mime": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "ext": "xlsx",
    },
    "xlsx": {
        "func": generate_xlsx,
        "mime": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "ext": "xlsx",
    },
    "pptx": {
        "func": generate_pptx,
        "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "ext": "pptx",
    },
}


async def generate_file(file_type: str, title: str, content: str) -> Tuple[str, str]:
    """
    Gera arquivo e faz upload ao Supabase.

    Returns:
        (url_download, mensagem)
    """
    from backend.core.config import get_config
    from backend.core.supabase_client import upload_to_supabase

    config = get_config()
    file_type = file_type.lower()

    if file_type not in FILE_GENERATORS:
        raise ValueError(f"Tipo de arquivo inválido: {file_type}. Suportados: {list(FILE_GENERATORS.keys())}")

    gen = FILE_GENERATORS[file_type]

    # Gerar arquivo
    file_bytes = gen["func"](title, content)

    # Filename seguro
    import re
    safe_title = re.sub(r'[^a-z0-9]', '_', title.lower())[:50]
    filename = f"{safe_title}.{gen['ext']}"

    # Upload ao Supabase
    url = upload_to_supabase(
        bucket=config.supabase_storage_bucket,
        filename=filename,
        file_content=file_bytes,
        content_type=gen["mime"],
    )

    return url, f"{file_type.upper()} gerado com sucesso."
